"""Módulo principal para manipulação do template PowerPoint do Diagnóstico Tributário.

Abre o template, atualiza os slides editáveis com dados do cliente e gráficos
gerados, e salva a apresentação final.
"""

import logging
import os
import shutil
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import (
    TEMPLATE_PATH,
    EDITABLE_SLIDES,
    SHAPE_MAP,
    CORES,
    FONTES,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Converte hex string (#RRGGBB) para RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


VERMELHO = _hex_to_rgb(CORES["vermelho"])
DOURADO = _hex_to_rgb(CORES["dourado"])
BRANCO = _hex_to_rgb(CORES["branco"])


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def _replace_text_in_shape(shape, new_text: str) -> None:
    """Substitui o texto de um shape preservando a formatação do primeiro run.

    Estratégia:
    - Coloca todo o novo texto no primeiro run do primeiro parágrafo.
    - Limpa o texto dos demais runs (mantém o XML intacto para preservar
      espaçamento / bullet se existirem).
    - Se houver múltiplos parágrafos, limpa todos exceto o primeiro.
    """
    if not shape.has_text_frame:
        logger.warning("Shape '%s' não tem text_frame; pulando.", shape.name)
        return

    tf = shape.text_frame
    paragraphs = tf.paragraphs

    if not paragraphs:
        return

    first_para = paragraphs[0]
    runs = first_para.runs

    if runs:
        # Preserva a formatação do primeiro run
        runs[0].text = str(new_text)
        # Limpa texto dos runs subsequentes
        for run in runs[1:]:
            run.text = ""
    else:
        # Sem runs: cria um novo run herdando o formato do parágrafo
        first_para.text = str(new_text)

    # Limpa parágrafos extras (mantendo o primeiro)
    for para in paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _replace_multiline_text(shape, new_text: str) -> None:
    """Substitui texto que pode conter '\\n', criando um parágrafo por linha.

    Preserva a formatação do primeiro run do primeiro parágrafo existente
    e a replica para as novas linhas.
    """
    if not shape.has_text_frame:
        logger.warning("Shape '%s' não tem text_frame; pulando.", shape.name)
        return

    tf = shape.text_frame
    lines = str(new_text).split("\n")

    # Captura referência de formatação do primeiro run
    ref_font = None
    ref_alignment = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        ref_run = tf.paragraphs[0].runs[0]
        ref_font = ref_run.font
        ref_alignment = tf.paragraphs[0].alignment

    # Limpa todo o text frame
    tf.clear()

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        run = para.add_run()
        run.text = line

        if ref_alignment is not None:
            para.alignment = ref_alignment

        # Copia propriedades de fonte se disponíveis
        if ref_font is not None:
            _copy_font_props(ref_font, run.font)


def _copy_font_props(src_font, dst_font) -> None:
    """Copia propriedades básicas de fonte de src para dst."""
    try:
        if src_font.size is not None:
            dst_font.size = src_font.size
        if src_font.bold is not None:
            dst_font.bold = src_font.bold
        if src_font.italic is not None:
            dst_font.italic = src_font.italic
        if src_font.name is not None:
            dst_font.name = src_font.name
        if src_font.color and src_font.color.rgb is not None:
            dst_font.color.rgb = src_font.color.rgb
    except Exception:
        # Algumas propriedades podem não estar definidas; ignorar
        pass


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def _replace_image(slide, shape_name: str, image_path: str) -> None:
    """Substitui um Picture shape por uma nova imagem, mantendo posição e tamanho.

    Remove o shape antigo e adiciona a nova imagem nas mesmas coordenadas.
    """
    target_shape = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target_shape = shape
            break

    if target_shape is None:
        logger.warning(
            "Shape de imagem '%s' não encontrado no slide %s; pulando.",
            shape_name,
            slide.slide_id,
        )
        return

    if not os.path.isfile(image_path):
        logger.error("Arquivo de imagem não encontrado: %s", image_path)
        return

    # Salva posição e tamanho
    left = target_shape.left
    top = target_shape.top
    width = target_shape.width
    height = target_shape.height

    # Remove o shape antigo
    sp = target_shape._element
    sp.getparent().remove(sp)

    # Adiciona nova imagem
    slide.shapes.add_picture(image_path, left, top, width, height)
    logger.info("Imagem '%s' substituída com sucesso.", shape_name)


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def _update_table_cell(cell, text: str, color=None) -> None:
    """Atualiza o texto de uma célula de tabela preservando a formatação.

    Args:
        cell: Objeto de célula do python-pptx.
        text: Novo texto.
        color: Se fornecido, aplica essa cor ao texto.
    """
    paragraphs = cell.text_frame.paragraphs
    if not paragraphs:
        return

    first_para = paragraphs[0]
    runs = first_para.runs

    if runs:
        runs[0].text = str(text)
        if color is not None:
            runs[0].font.color.rgb = color
        for run in runs[1:]:
            run.text = ""
    else:
        first_para.text = str(text)
        if color is not None and first_para.runs:
            first_para.runs[0].font.color.rgb = color

    # Limpa parágrafos extras
    for para in paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _find_table(slide):
    """Retorna o primeiro shape de tabela encontrado no slide, ou None."""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


# ---------------------------------------------------------------------------
# Slide-specific updaters
# ---------------------------------------------------------------------------

def _update_cover(slide, data: dict) -> None:
    """Slide 1 (idx 0) - Capa: atualiza nome do cliente."""
    mapping = SHAPE_MAP.get(0, {})
    dados_gerais = data.get("dados_gerais", {})
    logger.info("_update_cover: dados_gerais keys = %s", list(dados_gerais.keys()))
    logger.info("_update_cover: nome_cliente = '%s'", dados_gerais.get("nome_cliente", "<AUSENTE>"))

    for shape in slide.shapes:
        if shape.name in mapping:
            field = mapping[shape.name]
            value = dados_gerais.get(field, "")
            logger.info("_update_cover: shape='%s', field='%s', value='%s'", shape.name, field, value)
            if value:
                _replace_text_in_shape(shape, value)
                logger.info("Capa: '%s' → '%s'", shape.name, value)
            else:
                logger.warning("Capa: campo '%s' vazio, shape '%s' não atualizado.", field, shape.name)


def _update_indicators(slide, data: dict, charts: dict) -> None:
    """Slide 3 (idx 2) - Indicadores com gráfico donut."""
    mapping = SHAPE_MAP.get(2, {})
    indicadores = data.get("indicadores_resumo", {})
    dados_gerais = data.get("dados_gerais", {})

    for shape in slide.shapes:
        if shape.name not in mapping:
            continue

        field = mapping[shape.name]

        if field == "__donut_chart__":
            chart_path = charts.get("donut_chart")
            if chart_path:
                _replace_image(slide, shape.name, chart_path)
        else:
            # Tenta indicadores_resumo primeiro, depois dados_gerais como fallback
            value = indicadores.get(field, "")
            if value == "":
                value = dados_gerais.get(field, "")
            if value != "":
                _replace_text_in_shape(shape, value)
                logger.info("Indicadores: '%s' → '%s'", shape.name, value)


def _update_summary_table(slide, data: dict) -> None:
    """Slide 5 (idx 4) - Tabela de resumo de indicadores.

    Atualiza colunas 1 (valor) e 2 (percentual) das linhas 1-6.
    """
    table = _find_table(slide)
    if table is None:
        logger.warning("Tabela não encontrada no slide de resumo (idx 4).")
        return

    resumo = data.get("indicadores_resumo", {})

    # Mapeamento de linhas → chaves nos dados
    row_keys = {
        1: ("receita_valor", "receita_pct"),
        2: ("custo_variavel_valor", "custo_variavel_pct"),
        3: ("custo_fixo_valor", "custo_fixo_pct"),
        4: ("imposto_lucro_valor", "imposto_lucro_pct"),
        5: ("lucro_valor", "lucro_pct"),
        6: ("total_valor", "total_pct"),
    }

    for row_idx, (val_key, pct_key) in row_keys.items():
        if row_idx >= len(table.rows):
            logger.warning("Tabela de resumo tem menos linhas que o esperado.")
            break

        row = table.rows[row_idx]

        val = resumo.get(val_key, "")
        pct = resumo.get(pct_key, "")

        if len(row.cells) > 1 and val != "":
            _update_table_cell(row.cells[1], val)
        if len(row.cells) > 2 and pct != "":
            _update_table_cell(row.cells[2], pct)

    logger.info("Tabela de resumo atualizada.")


def _update_scenario_slide(slide, scenario_data: dict, slide_idx: int) -> None:
    """Slides 10/12/14 (idx 9/11/13) - Cenários com gauges.

    Atualiza text boxes de porcentagem LR, LP e texto de diferença.
    """
    mapping = SHAPE_MAP.get(slide_idx, {})

    for shape in slide.shapes:
        if shape.name not in mapping:
            continue

        field = mapping[shape.name]
        value = scenario_data.get(field, "")
        if value != "":
            _replace_text_in_shape(shape, value)
            logger.info("Cenário (idx %d): '%s' → '%s'", slide_idx, shape.name, value)


def _update_fiscal_benefits(slide, data: dict) -> None:
    """Slide 16 (idx 15) - Benefícios Fiscais (mapa do Brasil).

    Procura text boxes que começam com sigla de estado (ex: "PE - 11% de icms")
    e atualiza apenas esses shapes. Shapes de título/labels não são alterados.
    """
    beneficios = data.get("beneficios_fiscais", [])
    if not beneficios:
        logger.info("Sem dados de benefícios fiscais para atualizar.")
        return

    # Constrói lookup: sigla do estado → dados
    state_lookup = {}
    for item in beneficios:
        estado = item.get("estado", "").strip().upper()
        if estado:
            sigla = estado[:2]
            state_lookup[sigla] = item

    # Lista de siglas válidas de estados brasileiros para match seguro
    _UF_VALIDAS = {
        "AC", "AL", "AP", "AM", "BA", "CE", "DF", "ES", "GO", "MA",
        "MT", "MS", "MG", "PA", "PB", "PR", "PE", "PI", "RJ", "RN",
        "RS", "RO", "RR", "SC", "SP", "SE", "TO",
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text.strip()
        text_upper = text.upper()

        # Match seguro: o texto deve COMEÇAR com a sigla do estado
        # seguida de " -" ou " " (ex: "PE - 11% de icms", "SC - 2 a 3%...")
        matched_sigla = None
        for sigla in state_lookup:
            if sigla in _UF_VALIDAS and (
                text_upper.startswith(f"{sigla} -") or
                text_upper.startswith(f"{sigla} ") and len(text_upper) > 3
            ):
                matched_sigla = sigla
                break

        if matched_sigla is None:
            continue

        state_data = state_lookup[matched_sigla]
        icms_pct = state_data.get("icms_pct", "")
        lr_aliq1 = state_data.get("lr_aliq1", "")
        lr_aliq2 = state_data.get("lr_aliq2", "")
        lp_aliq1 = state_data.get("lp_aliq1", "")
        lp_aliq2 = state_data.get("lp_aliq2", "")

        new_lines = []
        if icms_pct:
            new_lines.append(f"{matched_sigla} - {icms_pct} de icms")
        else:
            new_lines.append(f"{matched_sigla}")
        if lr_aliq1 or lr_aliq2:
            new_lines.append(f"LR {lr_aliq1}/{lr_aliq2}")
        if lp_aliq1 or lp_aliq2:
            new_lines.append(f"LP {lp_aliq1}/{lp_aliq2}")

        if new_lines:
            _replace_multiline_text(shape, "\n".join(new_lines))
            logger.info("Benefício fiscal '%s' atualizado.", matched_sigla)


def _update_scenarios_table(slide, data: dict) -> None:
    """Slide 18 (idx 17) - Tabela de resumo de cenários (9x7).

    Atualiza dados nas linhas 1-8, aplicando vermelho em células marcadas.
    """
    table = _find_table(slide)
    if table is None:
        logger.warning("Tabela não encontrada no slide de cenários (idx 17).")
        return

    resumo = data.get("resumo_cenarios", [])

    for row_idx, row_data in enumerate(resumo, start=1):
        if row_idx >= len(table.rows):
            logger.warning("Tabela de cenários: dados excedem linhas disponíveis.")
            break

        row = table.rows[row_idx]

        # row_data deve ter as colunas na ordem da tabela
        col_keys = row_data.get("_col_keys", [])
        destaque = row_data.get("destaque_vermelho", [])

        if col_keys:
            for col_idx, key in enumerate(col_keys):
                if col_idx >= len(row.cells):
                    break
                val = row_data.get(key, "")
                # Substitui DIFAL vazio
                if "difal" in key.lower() and (val == "" or val is None):
                    val = " - "
                color = VERMELHO if key in destaque else None
                _update_table_cell(row.cells[col_idx], str(val), color=color)
        else:
            # Fallback: itera valores na ordem do dict (exceto metadados)
            meta_keys = {"_col_keys", "destaque_vermelho"}
            values = [v for k, v in row_data.items() if k not in meta_keys]
            for col_idx, val in enumerate(values):
                if col_idx >= len(row.cells):
                    break
                is_destaque = isinstance(destaque, list) and col_idx in destaque
                color = VERMELHO if is_destaque else None
                _update_table_cell(row.cells[col_idx], str(val), color=color)

    logger.info("Tabela de cenários atualizada.")


def _update_passivos(slide, data: dict) -> None:
    """Slide 19 (idx 18) - Gestão de Passivos.

    Procura text boxes com 'Federal' e 'Estadual' e substitui conteúdo.
    """
    passivos = data.get("gestao_passivos", {})
    federal_text = passivos.get("federal", "")
    estadual_text = passivos.get("estadual", "")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        current_text = shape.text_frame.text.strip().lower()

        # Caso especial: shape contém AMBOS "federal" e "estadual"
        if "federal" in current_text and "estadual" in current_text:
            combined_parts = []
            if federal_text:
                combined_parts.append(federal_text)
            if estadual_text:
                combined_parts.append(estadual_text)
            if combined_parts:
                _replace_multiline_text(shape, "\n\n".join(combined_parts))
                logger.info("Passivos Federal + Estadual atualizados (shape único).")
        elif "federal" in current_text and federal_text:
            _replace_multiline_text(shape, federal_text)
            logger.info("Passivos Federal atualizado.")
        elif "estadual" in current_text and estadual_text:
            _replace_multiline_text(shape, estadual_text)
            logger.info("Passivos Estadual atualizado.")


def _update_teses(slide, data: dict) -> None:
    """Slide 20 (idx 19) - Teses Tributárias.

    Limpa text boxes existentes e cria novos dinamicamente para cada tese.
    Cada tese tem: nome em branco bold + 'ECONOMIA: R$ X - PERÍODO' em dourado.
    """
    teses = data.get("teses_tributarias", [])
    if not teses:
        logger.info("Sem teses tributárias para exibir.")
        return

    # Remove text boxes existentes de teses (mantém shapes de fundo/título)
    # Shapes protegidos: TextBox 9 (subtítulo), TextBox 10 (título "Teses Tributárias")
    _PROTECTED_SHAPES = {"TextBox 9", "TextBox 10"}
    shapes_to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.name in _PROTECTED_SHAPES:
            continue
        text = shape.text_frame.text.strip().upper()
        # Heurística: remove shapes que parecem ser teses anteriores
        # (contêm "ECONOMIA" ou são text boxes editáveis na área de conteúdo)
        if "ECONOMIA" in text or "TESE" in text:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Cria novos text boxes para cada tese
    start_top = Inches(3.0)
    item_height = Inches(0.9)
    left = Inches(1.0)
    width = Inches(11.0)

    for i, tese in enumerate(teses):
        nome = tese.get("nome", tese.get("tese", f"Tese {i + 1}"))
        economia = tese.get("economia", "")
        periodo = tese.get("periodo", "")

        top = start_top + Emu(int(item_height * i))

        txBox = slide.shapes.add_textbox(left, top, width, item_height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Primeira linha: nome da tese
        para_nome = tf.paragraphs[0]
        run_nome = para_nome.add_run()
        run_nome.text = str(nome)
        run_nome.font.name = FONTES.get("titulo", "Montserrat")
        run_nome.font.size = Pt(16)
        run_nome.font.bold = True
        run_nome.font.color.rgb = BRANCO

        # Segunda linha: economia
        if economia:
            economia_text = f"ECONOMIA: {economia}"
            if periodo:
                economia_text += f" - {periodo}"

            para_eco = tf.add_paragraph()
            run_eco = para_eco.add_run()
            run_eco.text = economia_text
            run_eco.font.name = FONTES.get("titulo", "Montserrat")
            run_eco.font.size = Pt(13)
            run_eco.font.bold = True
            run_eco.font.color.rgb = DOURADO

    logger.info("%d tese(s) tributária(s) adicionada(s).", len(teses))


def _update_recuperacao_table(slide, data: dict) -> None:
    """Slide 21 (idx 20) - Tabela de Recuperação Tributária.

    Atualiza tabela existente e o VALOR TOTAL abaixo.
    """
    table = _find_table(slide)
    if table is None:
        logger.warning("Tabela não encontrada no slide de recuperação (idx 20).")
        return

    recuperacao = data.get("recuperacao_tributaria", [])
    total = data.get("recuperacao_tributaria_total", "R$ 0,00")

    for row_idx, row_data in enumerate(recuperacao, start=1):
        if row_idx >= len(table.rows):
            logger.warning("Recuperação: dados excedem linhas da tabela.")
            break

        row = table.rows[row_idx]

        # Espera-se colunas: RCT, IMPOSTO, VALOR CRÉDITO
        meta_keys = {"_col_keys", "destaque_vermelho"}
        values = [v for k, v in row_data.items() if k not in meta_keys]
        for col_idx, val in enumerate(values):
            if col_idx >= len(row.cells):
                break
            _update_table_cell(row.cells[col_idx], str(val))

    # Atualiza VALOR TOTAL em text box abaixo da tabela
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "VALOR TOTAL" in shape.text_frame.text.upper():
            _replace_text_in_shape(shape, f"VALOR TOTAL: {total}")
            logger.info("Recuperação: VALOR TOTAL → %s", total)
            break

    logger.info("Tabela de recuperação tributária atualizada.")


def _update_reforma_chart(slide, data: dict, charts: dict) -> None:
    """Slide 24 (idx 23) - Reforma Tributária.

    NÃO substitui gráfico de barras — o template PPTX já tem as barras desenhadas.
    Apenas atualiza os TextBoxes de valores acima das barras e as alíquotas.
    """
    mapping = SHAPE_MAP.get(23, {})
    reforma = data.get("reforma_tributaria", [])

    for shape in slide.shapes:
        if shape.name not in mapping:
            continue

        field = mapping[shape.name]

        if field == "aliquotas_texto":
            aliquotas = data.get("indicadores_resumo", {}).get("aliquotas_texto", "")
            if not aliquotas:
                # Tenta construir a partir de dados da reforma
                cbs = data.get("indicadores_resumo", {}).get("cbs_pct", "")
                ibs = data.get("indicadores_resumo", {}).get("ibs_pct", "")
                is_pct = data.get("indicadores_resumo", {}).get("is_pct", "")
                if cbs or ibs or is_pct:
                    aliquotas = f"Alíquotas:\nCBS: {cbs} IBS: {ibs} IS: {is_pct}"
            if aliquotas:
                _replace_multiline_text(shape, aliquotas)
                logger.info("Reforma: alíquotas atualizadas.")

    # Atualiza text boxes de valores das barras (TextBox 4-9)
    # TextBox 14 e 15 são duplicatas visuais de TextBox 4 e 5 (mesma posição, outra camada)
    bar_values = []
    for item in reforma:
        val = item.get("valor", item.get("valor_barra", ""))
        if val:
            bar_values.append(str(val))

    # Mapeamento: TextBox → índice no bar_values
    bar_textbox_map = {
        "TextBox 4": 0, "TextBox 5": 1, "TextBox 6": 2,
        "TextBox 7": 3, "TextBox 8": 4, "TextBox 9": 5,
        "TextBox 14": 0, "TextBox 15": 1,  # duplicatas de TextBox 4 e 5
    }
    for shape in slide.shapes:
        if shape.name in bar_textbox_map:
            idx = bar_textbox_map[shape.name]
            if idx < len(bar_values):
                _replace_text_in_shape(shape, bar_values[idx])
                logger.info("Reforma barra: '%s' → '%s'", shape.name, bar_values[idx])


def _update_sintese(slide, data: dict) -> None:
    """Slide 26 (idx 25) - Síntese do Diagnóstico.

    Usa mapeamento direto de shapes por nome para substituir conteúdo.
    """
    sintese = data.get("sintese_diagnostico", {})
    if not sintese:
        logger.info("Sem dados de síntese do diagnóstico.")
        return

    mapping = SHAPE_MAP.get(25, {})

    for shape in slide.shapes:
        if shape.name not in mapping:
            continue

        field = mapping[shape.name]
        value = sintese.get(field, "")
        if value:
            _replace_multiline_text(shape, str(value))
            logger.info("Síntese: '%s' → campo '%s'", shape.name, field)


# ---------------------------------------------------------------------------
# Slide dispatcher
# ---------------------------------------------------------------------------

# Maps slide index → updater function signature
# Each function receives (slide, data, charts) but some ignore charts
_SLIDE_UPDATERS = {
    0: lambda s, d, c: _update_cover(s, d),
    2: lambda s, d, c: _update_indicators(s, d, c),
    4: lambda s, d, c: _update_summary_table(s, d),
    9: lambda s, d, c: _update_scenario_slide(s, d.get("cenarios_comparativos", [{}])[0] if d.get("cenarios_comparativos") else {}, 9),
    11: lambda s, d, c: _update_scenario_slide(s, d.get("cenarios_comparativos", [{}])[1] if len(d.get("cenarios_comparativos", [])) > 1 else {}, 11),
    13: lambda s, d, c: _update_scenario_slide(s, d.get("cenarios_comparativos", [{}])[2] if len(d.get("cenarios_comparativos", [])) > 2 else {}, 13),
    15: lambda s, d, c: _update_fiscal_benefits(s, d),
    17: lambda s, d, c: _update_scenarios_table(s, d),
    18: lambda s, d, c: _update_passivos(s, d),
    19: lambda s, d, c: _update_teses(s, d),
    20: lambda s, d, c: _update_recuperacao_table(s, d),
    23: lambda s, d, c: _update_reforma_chart(s, d, c),
    25: lambda s, d, c: _update_sintese(s, d),
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def update_presentation(data: dict, charts: dict, output_path: str) -> str:
    """Gera a apresentação final a partir do template e dados do cliente.

    Args:
        data: Dados processados do cliente (retorno de data_processor.process_data).
        charts: Dict de nome_do_gráfico → caminho do arquivo PNG.
                Chaves esperadas: 'donut_chart', 'bar_chart'.
        output_path: Caminho completo para salvar o .pptx gerado.

    Returns:
        Caminho do arquivo gerado (mesmo que output_path).

    Raises:
        FileNotFoundError: Se o template não existir.
        Exception: Erros do python-pptx durante manipulação.
    """
    if not os.path.isfile(TEMPLATE_PATH):
        raise FileNotFoundError(f"Template não encontrado: {TEMPLATE_PATH}")

    # Garante que o diretório de saída exista
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    # Copia template para output (trabalha na cópia)
    shutil.copy2(TEMPLATE_PATH, output_path)
    logger.info("Template copiado para: %s", output_path)

    # Abre a cópia
    prs = Presentation(output_path)

    total_slides = len(prs.slides)
    logger.info("Apresentação aberta com %d slides.", total_slides)

    # Itera apenas os slides editáveis
    for slide_idx in EDITABLE_SLIDES:
        if slide_idx >= total_slides:
            logger.warning(
                "Slide idx %d fora do range (%d slides); pulando.",
                slide_idx,
                total_slides,
            )
            continue

        slide = prs.slides[slide_idx]
        updater = _SLIDE_UPDATERS.get(slide_idx)

        if updater is None:
            logger.warning("Sem updater para slide idx %d; pulando.", slide_idx)
            continue

        try:
            updater(slide, data, charts)
            logger.info("Slide idx %d atualizado com sucesso.", slide_idx)
        except Exception:
            logger.exception("Erro ao atualizar slide idx %d.", slide_idx)

    # Salva
    prs.save(output_path)
    logger.info("Apresentação salva em: %s", output_path)

    return output_path
